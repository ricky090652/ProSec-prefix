"""依 ProSec進度報告.pdf 的版面重建投影片，並補上 DPO full train 結果與 future work。

用法：.venv/bin/python docs/build_deck.py
輸出：docs/ProSec進度報告_v2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

BLACK = RGBColor(0, 0, 0)
RED = RGBColor(0xC0, 0x00, 0x00)
BLUE = RGBColor(0x44, 0x72, 0xC4)
GREY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1F, 0x7A, 0x3D)

LATIN = "Calibri"
EA = "Microsoft JhengHei"
MONO = "Consolas"


def style(run, size=14, bold=False, color=BLACK, latin=LATIN, ea=EA, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", ea if tag == "a:ea" else latin)
    return run


def textbox(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tf


def title(slide, text, size=36):
    tf = textbox(slide, 0.55, 0.28, 12.2, 0.95)
    style(tf.paragraphs[0].add_run(), size)
    tf.paragraphs[0].runs[0].text = text
    style(tf.paragraphs[0].runs[0], size, latin="Calibri Light")
    return tf


def bullets(slide, x, y, w, h, items, size=15, gap=4):
    """items: [(text, level, color, bold), ...] or plain strings"""
    tf = textbox(slide, x, y, w, h)
    first = True
    for it in items:
        if isinstance(it, str):
            it = (it, 0, BLACK, False)
        text, lvl, color, bold = (list(it) + [0, BLACK, False])[:4]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + text if text else ""
        style(r, size - lvl, bold=bold, color=color)
    return tf


def plain(slide, x, y, w, h, lines, size=14, mono=False, color=BLACK, gap=2):
    tf = textbox(slide, x, y, w, h)
    first = True
    for ln in lines:
        c = color
        b = False
        if isinstance(ln, tuple):
            ln, c = ln[0], ln[1]
            if len(ln) and isinstance(c, bool):
                ln, b, c = ln, c, color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ln
        style(r, size, bold=b, color=c, latin=MONO if mono else LATIN,
              ea=MONO if mono else EA)
    return tf


def table(slide, x, y, w, rows, col_w=None, size=12, head=True, row_h=0.32):
    nr, nc = len(rows), len(rows[0])
    shp = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w),
                                 Inches(row_h * nr))
    tbl = shp.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / total))
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci, cell in enumerate(row):
            txt, color, bold = cell if isinstance(cell, tuple) else (cell, BLACK, False)
            c = tbl.cell(ri, ci)
            c.margin_left = c.margin_right = Inches(0.06)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(txt)
            style(r, size, bold=bold or (head and ri == 0), color=color)
    return tbl


def box(slide, x, y, w, h, lines, outline=BLUE, size=13, rounded=True, fill=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = outline
    shape.line.width = Pt(1.25)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for ln in lines:
        txt, color, bold = (ln if isinstance(ln, tuple) else (ln, BLACK, False))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = txt
        style(r, size, bold=bold, color=color)
    return shape


def arrow(slide, x1, y1, x2, y2, color=BLUE):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(1.5)
    ln = c.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    return c


def label(slide, x, y, w, text, size=11, color=BLACK, align=PP_ALIGN.CENTER):
    tf = textbox(slide, x, y, w, 0.3)
    tf.paragraphs[0].alignment = align
    style(tf.paragraphs[0].add_run(), size, color=color)
    tf.paragraphs[0].runs[0].text = text
    return tf


def code(slide, x, y, w, h, lines, size=11.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    shape.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for ln in lines:
        txt, color = ln if isinstance(ln, tuple) else (ln, BLACK)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = txt
        style(r, size, color=color, latin=MONO, ea=MONO)
    return shape


def placeholder(slide, x, y, w, h, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = RGBColor(0xA6, 0xA6, 0xA6)
    shape.line.width = Pt(1)
    shape.line.dash_style = 4  # dashed
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style(p.add_run(), 13, color=GREY)
    p.runs[0].text = text
    return shape


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
S = lambda: prs.slides.add_slide(BLANK)

# ───────────────────────── 1 · Title ─────────────────────────
s = S()
tf = textbox(s, 0, 2.7, 13.333, 1.2)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
style(tf.paragraphs[0].add_run(), 54, latin="Calibri Light")
tf.paragraphs[0].runs[0].text = "進度報告"
tf2 = textbox(s, 0, 4.0, 13.333, 0.5)
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
style(tf2.paragraphs[0].add_run(), 20)
tf2.paragraphs[0].runs[0].text = "2026/09/03"

# ───────────────────────── 2 · Topic ─────────────────────────
s = S()
title(s, "Topic")
box(s, 0.6, 1.35, 5.6, 1.85, [
    ("SVEN (2023)", BLACK, True),
    ("• Prefix tuning for secure / insecure code generation", RED, False),
    ("• Freezes the base model", RED, False),
    ("• CodeGen (2022), torch 1.13 → not reusable", BLACK, False),
    ("• Real-world GitHub commits → small and noisy", BLACK, False),
], size=13)
box(s, 7.1, 1.35, 5.6, 1.85, [
    ("ProSec (2025)", BLACK, True),
    ("• Proactively synthesized vulnerability data", RED, False),
    ("• Preference optimization for fine-tuning", RED, False),
    ("• Adapts with LoRA", BLACK, False),
], size=13)
arrow(s, 3.4, 3.25, 5.4, 4.15)
arrow(s, 9.9, 3.25, 7.9, 4.15)
label(s, 2.1, 3.45, 2.2, "Prefix tuning", 12)
label(s, 9.9, 3.35, 3.2, "Data construct pipeline\npreference optimization", 12)
box(s, 2.9, 4.2, 7.5, 0.75, [
    ("ProSec's data and pipeline with ", BLACK, True),
], size=15)
# second run in red on the same paragraph
p = s.shapes[-1].text_frame.paragraphs[0]
style(p.add_run(), 15, bold=True, color=RED).text = "prefix (DPO) instead of LoRA"

# ─────────────────── 3 · ProSec data pipeline ───────────────────
s = S()
title(s, "Data syn & sel pipeline of ProSec")
placeholder(s, 0.6, 1.3, 7.4, 4.4, "貼上原本的 ProSec Figure 2")
bullets(s, 8.3, 1.35, 4.5, 4.2, [
    ("Step 1  Synthesize vulnerability-inducing instructions", 0, BLACK, True),
    ("50,000 instructions, clustered", 1, GREY, False),
    ("Step 2  Construct candidate preference data", 0, BLACK, True),
    ("Code LLM writes normal / fixed / vulnerable code", 1, GREY, False),
    ("Step 3  Select high-quality samples", 0, BLACK, True),
    ("Training-dynamics sampler + heuristic filter", 1, GREY, False),
], size=14)
box(s, 8.3, 4.55, 4.5, 1.15, [
    ("我們拿到的 45,785 筆已經是 Step 3 的輸出", RED, True),
    ("三項指紋：嚴格分區 · y_f 已去重 · D_norm 候選數卡在 2", BLACK, False),
    ("→ 再做一次 selection 是第二道過濾，不是復現", BLACK, False),
], size=12, outline=RED)

# ─────────────────── 4 · Prefix & LoRA ───────────────────
s = S()
title(s, "LoRA & Prefix")
label(s, 0.6, 1.25, 5.9, "LoRA — 在 transformer 內部層注入可訓練的 low-rank 矩陣", 14, align=PP_ALIGN.LEFT)
label(s, 7.0, 1.25, 5.9, "Prefix — 在每層的 key/value 序列前插入虛擬 token", 14, align=PP_ALIGN.LEFT)
placeholder(s, 0.6, 1.65, 5.9, 1.6, "貼上 LoRA 結構圖")
placeholder(s, 7.0, 1.65, 5.9, 1.6, "貼上 Prefix / KV 序列圖")
box(s, 0.6, 3.4, 5.9, 2.25, [
    ("Pros", BLACK, True),
    ("• 效能更接近全參數微調", BLACK, False),
    ("• 使用容易（開源解法多）", BLACK, False),
    ("Cons", BLACK, True),
    ("• 複雜任務調高 rank → 參數量與記憶體增加", BLACK, False),
    ("• 改變預訓練權重，切換需 merge / unmerge", BLACK, False),
], size=13)
box(s, 7.0, 3.4, 5.9, 2.25, [
    ("Pros", BLACK, True),
    ("• 參數量少，完整保留預訓練權重", BLACK, False),
    ("• 可依情境套用不同 prefix，切換≈換一組 KV", RED, False),
    ("• 同一個 batch 內可用不同 prefix（LoRA 做不到）", RED, False),
    ("Cons", BLACK, True),
    ("• 效能相對較差，且梯度只能經 softmax 回傳", BLACK, False),
], size=13)
table(s, 0.6, 5.8, 12.3, [
    ["", "可訓練參數", "adapter 大小 (bf16)", "實測最佳 lr"],
    ["LoRA r=8, all-linear", "12,582,912", "≈ 25.2 MB", "5e-6"],
    ["Prefix nvt=16", ("3,145,728", RED, True), ("≈ 6.3 MB", RED, True), ("5e-5", RED, True)],
], col_w=[3, 2, 2, 2], size=12.5)

# ─────────────── 5 · Different from paper ───────────────
s = S()
title(s, "Different from paper and follows")
table(s, 0.6, 1.4, 12.2, [
    ["", "ProSec paper", "This work", "Reason"],
    ["PEFT", "LoRA r=8, α=16", ("Prefix nvt = 16", RED, True), "Pluggable, flexible"],
    ["Objective", "SimPO (β=1.5, γ=0.5)", ("DPO (β=0.05)", RED, True),
     "SimPO 在全長訓練下兩臂都崩（prefix step ~450、LoRA step ~350）；論文 Table 8 自己驗證過 DPO"],
    ["Training length", "1500 steps @ batch 64", "800 steps @ batch 64", "≈ 1.1 epoch over 45,785"],
    ["LoRA target modules", "未指定", "all-linear", "我們自己的假設，論文只寫 r / α"],
    ["Evaluation",
     "PurpleLlama ∩ SafeCoder\n38 pairs / 694 cases（未釋出）\nHumanEval",
     ("Rebuilt 36 pairs / 693 cases\nHumanEval\nDegeneration test", RED, True),
     "子集必須自行重建\nThe security rate can be faked"],
], col_w=[2.2, 3, 3, 4], size=12, row_h=0.55)

# ─────────────── 6 · ProSec dataset ───────────────
s = S()
title(s, "ProSec dataset")
box(s, 0.5, 3.0, 2.9, 1.2, [
    ("Haiku-vul-inducing-", BLACK, True),
    ("instructions-clustered", BLACK, True),
    ("(50,000 instructions)", BLACK, False),
], size=13)
box(s, 4.1, 1.75, 3.4, 1.15, [
    ("Phi-3-mini generates", BLACK, True),
    ("Vulnerable + fixed code", BLUE, False),
    ("(validated by detector)", BLACK, False),
], size=12.5)
box(s, 4.1, 4.35, 3.4, 1.15, [
    ("Phi-3-mini generates", BLACK, True),
    ("Normal code + benign code", GREEN, False),
], size=12.5)
box(s, 8.1, 1.75, 3.2, 1.15, [
    ("D_sec, 27,400 pairs", BLUE, True),
    ("fixed ≻ vulnerable", BLACK, False),
], size=13)
box(s, 8.1, 4.35, 3.2, 1.15, [
    ("D_norm, 18,385 pairs", GREEN, True),
    ("normal ≻ over-secure", BLACK, False),
], size=13)
box(s, 11.7, 3.0, 1.5, 1.2, [("45,785", BLACK, True), ("pairs", BLACK, False)], size=14)
arrow(s, 3.4, 3.3, 4.05, 2.4)
arrow(s, 3.4, 3.9, 4.05, 4.8)
arrow(s, 7.55, 2.32, 8.05, 2.32)
arrow(s, 7.55, 4.92, 8.05, 4.92)
arrow(s, 11.35, 2.6, 11.9, 3.3)
arrow(s, 11.35, 4.6, 11.9, 3.9)

# ─────────────── 7 · Whole pipeline ───────────────
s = S()
title(s, "Whole pipeline")
label(s, 0.35, 1.35, 1.4, "training", 13, align=PP_ALIGN.LEFT)
box(s, 0.55, 1.75, 2.5, 0.95, [("ProSec dataset", BLACK, True), ("(45,785 pairs)", BLACK, False)])
box(s, 3.45, 1.75, 2.6, 0.95, [("(prompt, chosen,", BLACK, False), ("rejected) jsonl", BLACK, False)])
box(s, 6.45, 1.75, 3.0, 0.95, [("Prefix / LoRA training", BLACK, True),
                               ("(freeze phi-3-mini)", BLACK, False)],
    fill=RGBColor(0xE2, 0xEF, 0xDA))
box(s, 9.85, 1.75, 3.0, 0.95, [("Secure adapter", BLACK, True),
                               ("prefix ≈ 6.3 MB", RED, False)],
    fill=RGBColor(0xDE, 0xEA, 0xF6))
arrow(s, 3.1, 2.22, 3.4, 2.22)
arrow(s, 6.1, 2.22, 6.4, 2.22)
arrow(s, 9.5, 2.22, 9.8, 2.22)
c = slide = None
arrow(s, 11.35, 2.75, 11.35, 3.35, GREY)
arrow(s, 11.35, 3.35, 1.8, 3.35, GREY)
arrow(s, 1.8, 3.35, 1.8, 3.95, GREY)
label(s, 0.35, 3.55, 1.6, "evaluation", 13, align=PP_ALIGN.LEFT)
box(s, 0.55, 3.95, 2.7, 0.95, [("Generate response", BLACK, True),
                               ("693 × 10 samples × 2", BLACK, False)])
box(s, 3.65, 3.95, 2.7, 0.95, [("Normalize response", BLACK, True),
                               ("(wrap fence-less code)", BLACK, False)])
box(s, 6.75, 3.95, 2.7, 0.95, [("Static analyzer", BLACK, True),
                               ("semgrep / weggli / regex", BLACK, False)])
box(s, 9.85, 3.95, 2.7, 0.95, [("Score detect", BLACK, True),
                               ("(vulnerable ratio)", BLACK, False)])
arrow(s, 3.3, 4.42, 3.6, 4.42)
arrow(s, 6.4, 4.42, 6.7, 4.42)
arrow(s, 9.5, 4.42, 9.8, 4.42)
box(s, 1.9, 5.35, 4.2, 0.9, [("Utility test", BLACK, True),
                             ("HumanEval / MultiPL-E", BLACK, False)])
box(s, 7.0, 5.35, 4.2, 0.9, [("Check degeneration", BLACK, True),
                             ("length / parse rate / stub rate", BLACK, False)])

# ─────────────── 8 · Normalize response ───────────────
s = S()
title(s, "Normalize response")
bullets(s, 0.6, 1.15, 12.2, 1.15, [
    "SAST only detects code inside a fence（```）",
    ("Fence 是 Markdown 語法，沒有任何地方能設定——模型加不加是 instruction tuning 學來的習慣", 1, BLACK, False),
    ("而 benchmark 的 prompt 結尾寫著 \"Only return the code, don't include any other information\" → 模型不加 fence", 1, RED, False),
], size=14)
label(s, 0.6, 2.35, 6.0, "ProSec  detect_all.py", 12.5, align=PP_ALIGN.LEFT)
code(s, 0.6, 2.65, 6.0, 1.85, [
    "def parse_code_blocks(text):",
    "    code_blocks = []",
    "    in_code_block = False",
    "    for line in text.split('\\n'):",
    "        if '```' in line:",
    "            in_code_block = not in_code_block",
    "        elif in_code_block:",
    "            code_blocks.append(line)",
    ("    return '\\n'.join(code_blocks)   # 沒 fence → \"\"", RED),
], size=11)
label(s, 7.0, 2.35, 6.0, "eval/normalize_responses.py", 12.5, align=PP_ALIGN.LEFT)
code(s, 7.0, 2.65, 6.0, 1.05, [
    "if has_code_block(resp):",
    "    new.append(resp)                # 已有 fence，不動",
    "else:",
    ("    new.append(f\"```\\n{resp}\\n```\")   # 補上 fence", RED),
], size=11)
table(s, 7.0, 3.9, 6.0, [
    ["模型輸出", "抽取結果", "判定"],
    ["有 fence", "完整程式碼", "不安全 ✔"],
    [("沒有 fence", RED, True), ("\"\" 空字串", RED, True), ("「安全」← 錯", RED, True)],
    ["補 fence 後", "完整程式碼", "不安全 ✔"],
], col_w=[2, 2, 2], size=12)
box(s, 0.6, 4.75, 6.0, 1.35, [
    ("84% 的 code block 是空的", RED, True),
    ("漏洞率被稀釋約 6 倍，而整條鏈沒有任何一步報錯", BLACK, False),
    ("修完後 base 才與論文對齊：JS 52.24 vs 52.24", BLACK, False),
], size=13, outline=RED)

# ─────────────── 9 · Degeneration & additional check ───────────────
s = S()
title(s, "Degeneration & Additional check")
bullets(s, 0.6, 1.15, 12.2, 0.75, [
    "SAST can only report vulnerabilities in the code it can read",
    ("Anything that reduces readable code lowers the vulnerability rate", 1, BLACK, False),
], size=14)
table(s, 0.6, 2.0, 12.2, [
    ["Check", "Prevents", "Threshold"],
    ["Degeneration gate\n• code length\n• AST parse / bracket\n• stub / TODO / {}",
     "模型靠少寫程式碼來降低漏洞率",
     "長度保留 ≥ 90%\n語法完整率 Δ ≥ −3 pt\nstub 率 Δ ≤ +2 pt"],
    ["Pairwise drop",
     "生成被截斷時，若只丟單邊會破壞比較基礎",
     "任一邊截斷就兩邊一起丟"],
    ["Suppression ratio",
     "把「兩側一起壓低」誤判成模型學得更好",
     "Δrejected ÷ Δchosen；≈ 1 為失敗，> 2 健康"],
], col_w=[2.6, 4.6, 4.0], size=12, row_h=0.62)
label(s, 0.6, 4.95, 1.0, "EX：", 13, align=PP_ALIGN.LEFT)
table(s, 1.5, 4.9, 7.2, [
    ["max_new_tokens = 512", "漏洞率 Δ", "語法完整率 Δ", "程式碼長度"],
    ["LoRA ck800", ("−25.60", RED, True), ("−46.00", RED, True), "+75.6%"],
    ["Prefix ck800", "−4.40", "−4.60", "+1.6%"],
], col_w=[2.4, 1.6, 1.8, 1.6], size=12)
box(s, 9.1, 4.9, 3.8, 1.3, [
    ("程式碼被截斷會同時造成", BLACK, False),
    ("語法完整率下降 + 漏洞率下降", RED, True),
    ("提高到 1024 後：−46.00 → +0.25", BLACK, True),
], size=12.5)

# ─────────────── 10 · DPO lr sweep ───────────────
s = S()
title(s, "DPO lr sweep", 32)
label(s, 0.6, 1.05, 12.2, "β=0.05 · batch 64 · 16,000 samples = 250 steps · seed 42 · max_length 2048 · 兩臂同一張網格、同一條選擇規則",
      13, align=PP_ALIGN.LEFT)
label(s, 0.6, 1.45, 12.2, "選擇規則：比值 ≥ 3 且梯度未爆的最高 lr", 13, color=RED, align=PP_ALIGN.LEFT)
label(s, 0.6, 1.85, 3.0, "LoRA r=8 all-linear", 14, align=PP_ALIGN.LEFT)
table(s, 0.6, 2.15, 12.2, [
    ["lr", "chosen Δ", "rejected Δ", "壓rej/壓chosen", "margin early→late", "acc late", "grad 峰值", "chosen 保留/token"],
    [("5e-6 ✔", RED, True), "−0.0026 → −0.1500 (−0.147)", "−0.0127 → −0.6952 (−0.682)",
     ("4.63", RED, True), "0.010 → 0.545", "0.724", "58", ("99.3%", RED, True)],
    ["2e-5", "−0.0330 → −3.4503 (−3.417)", "−0.1687 → −8.2426 (−8.074)", "2.36", "0.136 → 4.792", "0.975", "106", "86.0%"],
    ["5e-5", "−0.1972 → −17.2034 (−17.006)", "−1.0295 → −34.0913 (−33.062)", "1.94", "0.832 → 16.888", "0.989", "366", "47.0%"],
    ["1e-4", "−1.0352 → −36.6601 (−35.625)", "−2.9977 → −69.9023 (−66.905)", "1.88", "1.962 → 33.242", "0.995", "569", "20.0%"],
    ["2e-4", "−3.2846 → −35.3859 (−32.101)", "−7.5961 → −83.7929 (−76.197)", "2.37", "4.311 → 48.407", "0.997", "608", "21.2%"],
], col_w=[0.8, 2.6, 2.6, 1.4, 1.7, 1.0, 1.0, 1.4], size=10.5, row_h=0.3)
label(s, 0.6, 4.15, 3.0, "Prefix nvt=16", 14, align=PP_ALIGN.LEFT)
table(s, 0.6, 4.45, 12.2, [
    ["lr", "chosen Δ", "rejected Δ", "壓rej/壓chosen", "margin early→late", "acc late", "grad 峰值", "chosen 保留/token"],
    ["5e-6", "−0.7780 → −0.7383 (+0.040)", "−0.8421 → −0.8575 (−0.015)", "∞", "0.064 → 0.119", "0.554", "45", "96.8%"],
    ["2e-5", "−0.7735 → −0.7394 (+0.034)", "−0.8484 → −0.9466 (−0.098)", "∞", "0.075 → 0.207", "0.612", "34", "96.8%"],
    [("5e-5 ✔", RED, True), "−0.7670 → −0.8250 (−0.058)", "−0.8613 → −1.2268 (−0.365)",
     ("6.30", RED, True), "0.094 → 0.402", "0.692", "67", ("96.4%", RED, True)],
    ["1e-4", "−0.7637 → −1.0707 (−0.307)", "−0.8921 → −1.6975 (−0.805)", "2.62", "0.129 → 0.627", "0.730", ("3722", RED, True), "95.4%"],
    ["2e-4", "−0.8065 → −4.2191 (−3.413)", "−1.0106 → −5.2037 (−4.193)", "1.23", "0.204 → 0.985", "0.647", ("3672", RED, True), "83.1%"],
], col_w=[0.8, 2.6, 2.6, 1.4, 1.7, 1.0, 1.0, 1.4], size=10.5, row_h=0.3)
label(s, 0.6, 6.45, 12.2,
      "兩個最佳點差 10 倍。同一個 5e-5：prefix 在最佳點（6.30），LoRA 已只剩 47% 機率保留（1.94）。",
      13, color=RED, align=PP_ALIGN.LEFT)

# ─────────────── 11 · DPO full train — training ───────────────
s = S()
title(s, "DPO full train — training dynamics")
label(s, 0.6, 1.15, 12.2, "β=0.05 · batch 64 · 800 steps · seed 42 · 兩臂唯一差別是 --peft_method",
      14, align=PP_ALIGN.LEFT)
table(s, 0.6, 1.75, 10.0, [
    ["", "LoRA  lr = 5e-6", "Prefix  lr = 5e-5"],
    ["可訓練參數", "12,582,912", ("3,145,728", RED, True)],
    ["chosen 每 token 機率保留", "88.8%", ("94.6%", RED, True)],
    ["rewards/accuracies", "0.975", "0.753"],
    ["rewards/margins", "4.11", "0.82"],
    ["grad_norm 峰值", "87", ("18,657", RED, True)],
], col_w=[3.4, 3.3, 3.3], size=14, row_h=0.42)
box(s, 0.6, 4.6, 12.2, 1.05, [
    ("prefix 的 18,657 不是發散：5 次尖峰、量級遞減，chosen 每次都回升（step 245 一度 −3.24，最終 −1.27）", BLACK, False),
    ("DPO 的 reference model 錨住了絕對機率——同樣的尖峰在 SimPO 下是不可逆的", BLACK, False),
], size=13)
box(s, 0.6, 5.85, 12.2, 0.7, [
    ("⚠ 兩臂的參數量差 4 倍，所以任何差距目前都是被混淆的", RED, True),
], size=14, outline=RED)

# ─────────────── 12 · Security result ───────────────
s = S()
title(s, "DPO full train — security")
label(s, 0.6, 1.15, 12.2,
      "100 prompts × 5 samples · C / C++ · max_new_tokens = 1024 · 成對移除截斷樣本 · 標準誤 ≈ 2.2 pt",
      13.5, align=PP_ALIGN.LEFT)
table(s, 0.6, 1.65, 9.6, [
    ["Vulnerable code ratio", "OFF", "ON", "Δ", "相對降幅"],
    ["LoRA", "49.41%", "31.62%", ("−17.80", RED, True), ("−36.0%", RED, True)],
    ["Prefix", "49.27%", "43.45%", ("−5.82", RED, True), "−11.8%"],
    ["Prefix 分語言", "—", "—", "c −6.58 / cpp −5.04", "—"],
], col_w=[3, 1.6, 1.6, 2, 1.8], size=13, row_h=0.4)
label(s, 0.6, 3.5, 6.0, "Degeneration gate", 14, align=PP_ALIGN.LEFT)
table(s, 0.6, 3.85, 12.2, [
    ["", "語法完整率 Δ\n(未截斷者)", "程式碼長度\nvs OFF", "截斷率", "守門判定"],
    ["LoRA", ("+0.25", RED, True), "226%", "14.4%", "通過"],
    ["Prefix", ("−4.74", RED, True), "103%", "2.6%", ("未過", RED, True)],
], col_w=[2, 2.6, 2.2, 1.8, 2.2], size=13, row_h=0.5)
box(s, 0.6, 5.55, 12.2, 1.15, [
    ("Prefix 這一臂沒有通過守門（門檻 −3），所以 −5.82 不得單獨回報，必須並列退化指標", RED, True),
    ("LoRA 通過守門，但程式碼長度 226%、截斷率 14.4%，是否在鑽偵測器漏洞仍未釐清", BLACK, False),
], size=13, outline=RED)

# ─────────────── 13 · vs paper ───────────────
s = S()
title(s, "Security — 與論文對照")
table(s, 0.6, 1.5, 12.2, [
    ["", "PEFT · Objective", "Vulnerable ratio", "相對降幅", "評測範圍"],
    ["ProSec Table 1", "LoRA · SimPO", "50.57 → 33.47", "−33.8%", "5 語言 / 694 題 / 10 samples"],
    ["ProSec Table 8", "LoRA · DPO", "40.76 → 34.65", "−15.0%", "隨機子集"],
    [("This work", RED, True), ("LoRA · DPO", RED, True), ("49.41 → 31.62", RED, True),
     ("−36.0%", RED, True), ("C / C++ / 100 題 / 5 samples", RED, True)],
    ["This work", "Prefix · DPO", "49.27 → 43.45", "−11.8%", "C / C++ / 100 題 / 5 samples"],
], col_w=[2.2, 2.4, 2.6, 1.8, 3.2], size=13, row_h=0.45)
box(s, 0.6, 4.2, 12.2, 1.5, [
    ("LoRA 復現的相對降幅與論文主結果同量級 → 資料前處理與訓練管線沒有問題", BLACK, True),
    ("", BLACK, False),
    ("但還不能說「對上了」：① 只跑 C / C++，正是 base 最脆弱的兩種語言（C 的 base 就有 ~70%）", BLACK, False),
    ("② 100 題不是 693 題　③ 我們用 DPO，論文主結果用 SimPO", BLACK, False),
], size=13)

# ─────────────── 14 · Utility ───────────────
s = S()
title(s, "DPO full train — utility (HumanEval)")
label(s, 0.6, 1.15, 12.2, "164 題 · greedy · max_new_tokens = 2048 · 截斷 0 / 164",
      13.5, align=PP_ALIGN.LEFT)
table(s, 0.6, 1.7, 8.5, [
    ["pass@1", "OFF (base)", "ON", "Δ"],
    ["LoRA", "70.73%", "66.46%", "−4.27"],
    ["Prefix", "70.73%", "59.76%", ("−10.98", RED, True)],
], col_w=[2.4, 2.2, 2.2, 2.0], size=14, row_h=0.42)
bullets(s, 0.6, 3.35, 12.2, 1.9, [
    ("H2 被推翻——原假設是 prefix 傷 utility 較少，實測是 LoRA 的 2.6 倍", 0, RED, True),
    ("加上安全性 −5.82 vs −17.80 → prefix 兩個軸都輸：學得比較少，卻傷得比較重", 0, RED, True),
    ("退化守門沒有預測到：prefix 長度保留 103%、看起來完好，pass@1 卻掉 11 pt", 0, BLACK, True),
    ("三條守門條件是「有沒有大幅退化」的篩子，不能代替 pass@1", 1, BLACK, False),
    ("逐題檢查：27 題退步 / 9 題進步，淨掉 18 題。抽查確認是真的能力退步，不是生成或抽取的 bug", 0, BLACK, False),
], size=13.5)
label(s, 0.6, 5.35, 6.0, "附帶發現：舊的 utility 數字被截斷污染", 14, color=RED, align=PP_ALIGN.LEFT)
table(s, 0.6, 5.7, 6.4, [
    ["max_new_tokens", "base HumanEval pass@1"],
    ["512（舊實驗用的）", "64.63%"],
    ["2048（現在）", ("70.73%", RED, True)],
], col_w=[3.2, 3.2], size=13, row_h=0.35)
box(s, 7.4, 5.7, 5.5, 1.05, [
    ("光是 base 就被低估 6.1 pt", RED, True),
    ("→ Exp1 / Exp2 / MultiPL-E 的舊 utility 數字全部不可信", BLACK, False),
], size=13)

# ─────────────── 15 · Current judgment ───────────────
s = S()
title(s, "目前的判定")
table(s, 0.6, 1.5, 12.2, [
    ["假設", "內容", "實測", "狀態"],
    ["H1", "Prefix 能以遠少的參數取得大部分安全性改善", "−5.82 vs LoRA −17.80（約 1/3）", ("被推翻", RED, True)],
    ["H2", "Prefix 的 utility 掉幅小於 LoRA", "−10.98 vs LoRA −4.27（2.6 倍）", ("被推翻", RED, True)],
], col_w=[1.2, 5.0, 3.8, 2.2], size=13.5, row_h=0.5)
box(s, 0.6, 3.4, 12.2, 1.35, [
    ("但兩個假設都還不能結案——目前的比較在參數量上差 4 倍", RED, True),
    ("Prefix nvt=16 是 3,145,728；LoRA all-linear 是 12,582,912", BLACK, False),
    ("「prefix 較弱」與「prefix 只拿到四分之一預算」目前分不開，只有前者叫能力邊界", BLACK, False),
], size=13.5, outline=RED)
label(s, 0.6, 5.0, 6.0, "已經排除的可能性", 14, align=PP_ALIGN.LEFT)
table(s, 0.6, 5.35, 12.2, [
    ["懷疑", "檢查方式", "結果"],
    ["截斷讓 LoRA 的安全性虛高", "成對移除截斷後重算", "排除：差距 12.6 → 12.0 pt"],
    ["截斷讓 pass@1 被低估", "max_new_tokens 提到 2048", "排除：截斷 0 / 164"],
    ["prefix 與 KV cache 的互動有 bug", "逐題檢查 27 題退步", "排除：是真的能力退步"],
    ["資料前處理或訓練管線有錯", "LoRA 走同一支程式，只差 peft config", "排除：LoRA 行為合理"],
], col_w=[4.0, 4.2, 4.0], size=12.5, row_h=0.36)

# ─────────────── 16 · Future work ───────────────
s = S()
title(s, "Future work")
label(s, 0.6, 1.2, 12.2, "① Parameter matching（最優先，做完才有資格對 H1 / H2 下結論）", 15,
      color=RED, align=PP_ALIGN.LEFT)
table(s, 0.9, 1.6, 10.5, [
    ["配對", "Prefix", "LoRA", "可訓練參數"],
    ["高容量", ("nvt = 64（待跑）", RED, True), "all-linear ✔ 已有", "12,582,912"],
    ["低容量", "nvt = 16 ✔ 已有", ("只掛 qkv_proj（待跑）", RED, True), "3,145,728"],
], col_w=[1.6, 3.2, 3.2, 2.5], size=13, row_h=0.36)
label(s, 0.6, 3.05, 12.2, "② Masked DPO — 只對真正不同的 token 算 loss", 15, align=PP_ALIGN.LEFT)
bullets(s, 0.9, 3.4, 12.0, 1.1, [
    ("動機：edit ratio ρ 中位數 0.358 → 64% 的 token 兩邊相同，對偏好訊號沒有貢獻卻完整計入 loss", 0, BLACK, False),
    ("外部驗證：PTC (DSN 2025) 用同樣的 diff mask，no control 55.0% → line 71.3% → mixed 86.7%", 0, BLACK, False),
    ("前置檢查（成敗關鍵）：17 個 edit hunk 有多少真的落在分析器標記的漏洞行上？若多為無關重寫則放棄", 0, RED, False),
], size=12.5)
label(s, 0.6, 4.75, 12.2, "③ CWE-specific 多 Prefix — 凸顯可插拔特性", 15, align=PP_ALIGN.LEFT)
bullets(s, 0.9, 5.1, 12.0, 1.1, [
    ("一個 base model 掛多組 ≈6 MB 模組，切換只是換一組 past_key_values", 0, BLACK, False),
    ("LoRA 做不到的是「同一個 batch 內不同請求用不同 policy」→ 要量的是切換延遲與混合 batch 吞吐量", 0, RED, False),
    ("必要對照：用同樣大小的隨機混合子集訓一個通用 prefix，否則分不出「專用有效」還是「資料變少」", 0, BLACK, False),
], size=12.5)
label(s, 0.6, 6.45, 12.2, "④ 補完評測：MultiPL-E 兩臂 · 全套 693 題 × 10 samples × 5 語言", 15,
      align=PP_ALIGN.LEFT)

out = "docs/ProSec進度報告_v2.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
